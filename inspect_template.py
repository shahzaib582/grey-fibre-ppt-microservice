import sys
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE


def get_title(slide):
    # Prefer the title placeholder
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and title_shape.has_text_frame:
        txt = title_shape.text_frame.text.strip()
        if txt:
            return txt

    # Fallback: first non-empty text frame
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            txt = shp.text_frame.text.strip()
            if txt:
                return txt
    return "(no title text)"

def inspect_template(path: str):
    prs = Presentation(path)
    print(f"Template: {path}")
    print(f"Total slides: {len(prs.slides)}")
    print()

    if len(prs.slides) < 7:
        print("Deck has fewer than 7 slides, cannot inspect slide 7.")
        return

    s7 = prs.slides[6]
    s5 = prs.slides[4]  # divider
    title7 = get_title(s7)
    layout7 = s7.slide_layout
    layout5 = s5.slide_layout
    print("=== Slide 5 (divider) vs Slide 7 (content) ===")
    print(f"Slide 5 layout: {getattr(layout5, 'name', '') or '(no name)'}, same as 7? {layout5 is layout7}")
    print(f"Slide 7 layout: {getattr(layout7, 'name', '') or '(no name)'}")
    print(f"Slide 5 follow_master_bg: {getattr(s5, 'follow_master_background', '?')}")
    print(f"Slide 7 follow_master_bg: {getattr(s7, 'follow_master_background', '?')}")
    print()

    # Show all layouts (there may be more than one)
    print("=== All slide layouts in deck ===")
    for i, layout in enumerate(prs.slide_layouts):
        name = getattr(layout, "name", "") or "(no name)"
        print(f"  Layout {i}: '{name}'")
    print()

    # Scan for Key Findings slides: divider (no chart) vs content (has chart)
    print("=== Slides with 'Key Findings' (divider vs content) ===")
    for i, s in enumerate(prs.slides):
        text = ""
        for shp in s.shapes:
            if getattr(shp, "has_text_frame", False):
                text += shp.text_frame.text + " "
        if "key findings" not in text.lower():
            continue
        has_chart = any(getattr(shp, "has_chart", False) for shp in s.shapes)
        title = get_title(s)
        print(f"  Slide {i+1}: title={title!r}, has_chart={has_chart} -> {'CONTENT (use this)' if has_chart else 'DIVIDER (skip)'}")
    print()

    # List ALL shapes on slide 7 to find where light grey + blue footer come from
    print("=== All shapes on slide 7 (content slide) ===")
    slide_height = prs.slide_height
    for shp in s7.shapes:
        top = getattr(shp, "top", None)
        left = getattr(shp, "left", None)
        height = getattr(shp, "height", None)
        kind = []
        if getattr(shp, "is_placeholder", False):
            kind.append("placeholder")
        if getattr(shp, "has_chart", False):
            kind.append("chart")
        if getattr(shp, "has_text_frame", False):
            kind.append("text")
        name = getattr(shp, "name", "(no name)")
        text = ""
        if getattr(shp, "has_text_frame", False):
            text = shp.text_frame.text.replace("\n", " ").strip()[:50]
        print(f"- {name!r}: top={top}, left={left}, h={height}, kind={','.join(kind) or 'shape'}, text={text!r}")

    # Dump cSld XML to see background structure
    print("\n=== Slide 7 cSld XML (first 2500 chars) ===")
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        el = s7._element
        cSld = el.find(qn("p:cSld"))
        if cSld is not None:
            xml_str = etree.tostring(cSld, encoding="unicode", pretty_print=True)
            print(xml_str[:2500])
        else:
            print("No p:cSld found")
    except Exception as e:
        print(f"Error: {e}")

def verify_pipeline_output(path: str):
    """Verify pipeline output: sections, transition titles, placeholders, slide numbers."""
    from survey_pipeline.utils import PLACEHOLDER, is_section_divider

    prs = Presentation(path)
    print(f"Verify: {path}")
    print(f"Total slides: {len(prs.slides)}")
    print()

    def _get_title(slide):
        title_types = (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if getattr(shape, "is_placeholder", False) and getattr(shape, "placeholder_format", None):
                ph_type = getattr(shape.placeholder_format, "type", None)
                if ph_type == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
                    continue
                if ph_type in title_types:
                    text = shape.text_frame.text.strip()
                    if text:
                        return text
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if getattr(shape, "is_placeholder", False) and getattr(shape.placeholder_format, "type", None) == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
                continue
            text = shape.text_frame.text.strip()
            if text:
                return text
        return ""

    # 1. Section dividers and transition slides
    print("=== Section dividers and transition slides ===")
    section_indices = []
    for idx, slide in enumerate(prs.slides):
        sec_name = is_section_divider(slide)
        if sec_name:
            section_indices.append((idx, sec_name))

    for j, (start_idx, section_name) in enumerate(section_indices):
        end_idx = section_indices[j + 1][0] if j + 1 < len(section_indices) else len(prs.slides)
        print(f"\nSection '{section_name}' (slides {start_idx+1}–{end_idx}):")
        for k in range(start_idx + 1, min(start_idx + 5, end_idx)):
            slide = prs.slides[k]
            title = _get_title(slide)
            has_sldnum = any(
                getattr(shp, "is_placeholder", False)
                and getattr(shp, "placeholder_format", None)
                and getattr(shp.placeholder_format, "type", None) == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER
                for shp in slide.shapes
            )
            print(f"  Slide {k+1}: title={title[:70]!r}  sldNum={has_sldnum}")

    # 2. Remaining placeholders
    print("\n=== Remaining placeholders ===")
    placeholder_slides = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and PLACEHOLDER in shape.text_frame.text:
                placeholder_slides.append(i + 1)
                if len(placeholder_slides) <= 3:
                    print(f"  Slide {i+1}: shape {getattr(shape, 'name', '?')!r}")
                break
    if not placeholder_slides:
        print("  [OK] None")
    else:
        print(f"  Slides with {PLACEHOLDER}: {placeholder_slides}")

    # 3. Slide number on key content slide (e.g. slide 7)
    print("\n=== Slide number placeholder on slide 7 ===")
    if len(prs.slides) >= 7:
        s7 = prs.slides[6]
        for shp in s7.shapes:
            if getattr(shp, "is_placeholder", False) and getattr(shp, "placeholder_format", None):
                if getattr(shp.placeholder_format, "type", None) == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
                    print(f"  Found: {getattr(shp, 'name', '?')!r}")
                    break
        else:
            print("  [MISSING] No slide number placeholder found")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python inspect_template.py TEMPLATE.pptx")
        print("       python inspect_template.py --verify OUTPUT.pptx")
        sys.exit(1)
    if sys.argv[1] == "--verify":
        if len(sys.argv) != 3:
            print("Usage: python inspect_template.py --verify OUTPUT.pptx")
            sys.exit(1)
        verify_pipeline_output(sys.argv[2])
    else:
        inspect_template(sys.argv[1])