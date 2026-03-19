"""
Shared utilities for the 3-Step Survey Slide Automation pipeline.
"""

import re
import copy
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PLACEHOLDER = "{Insert Finding Here}"


def _sanitize_pptx_text(text: str) -> str:
    """Remove control chars that can corrupt PPTX XML."""
    if not text or not isinstance(text, str):
        return text or ""
    return "".join(c for c in text if c == "\n" or c == "\r" or c == "\t" or ord(c) >= 32 or ord(c) == 0x2028 or ord(c) == 0x2029)


def _strip_question_ids(text: str) -> str:
    """Remove question IDs (Q1, Q7, etc.) from LLM output—readers must not need to look at question slides."""
    if not text or not isinstance(text, str):
        return text
    # "In Q7, " or "in Q7, " at start
    text = re.sub(r"\b[Ii]n\s+Q\d+,\s*", "", text, flags=re.IGNORECASE)
    # "Q7, " or "Q7: " anywhere
    text = re.sub(r"\bQ\d+\s*[,:\s]+\s*", " ", text, flags=re.IGNORECASE)
    # " (Q7)" or "(Q7)"
    text = re.sub(r"\s*\(Q\d+\)\s*", " ", text, flags=re.IGNORECASE)
    # "Q7 " at start of line/bullet
    text = re.sub(r"(^|\n)\s*Q\d+\s+", r"\1", text, flags=re.IGNORECASE)
    return re.sub(r"  +", " ", text).strip()


# ── Section divider names (extend for known templates; new sections auto-detected) ──
SECTION_NAMES = [
    "Mood",
    "Favorability",
    "Ballot",
    "Positioning",
    "Pro Gill Messages",
    "Anti Gill Messages",
    "Anti Malinowski Messages",
    "Demographics",
]

# Cached style for key-finding text (from the placeholder in the template)
KEY_FINDING_STYLE: dict | None = None


# ═══════════════════════════════════════════════════════════
#  Question spec parsing
# ═══════════════════════════════════════════════════════════

def parse_question_spec(slide_text: str):
    """
    Detect question specification from slide text.
    Returns:
        ("single", int)         — e.g. Question 5:
        ("range", (int, int))   — e.g. Questions 6-16:
        None                    — no question detected
    """
    # Try range pattern first (Questions X-Y:)
    m = re.search(r"\bQuestions?\s+(\d+)\s*[-–—]\s*(\d+)\s*:", slide_text, flags=re.IGNORECASE)
    if m:
        a, b = int(m.group(1)), int(m.group(2))
        if a > b:
            a, b = b, a
        return ("range", (a, b))

    # Single question pattern (Question N:)
    m = re.search(r"\bQuestion\s+(\d+)\s*:", slide_text, flags=re.IGNORECASE)
    if m:
        return ("single", int(m.group(1)))

    return None


def get_question_ids(qspec):
    """Convert a question spec to a list of question IDs (e.g. ['Q6', 'Q7', ...])."""
    if qspec is None:
        return []
    if qspec[0] == "single":
        return [f"Q{qspec[1]}"]
    elif qspec[0] == "range":
        a, b = qspec[1]
        return [f"Q{i}" for i in range(a, b + 1)]
    return []


# ═══════════════════════════════════════════════════════════
#  Data selection
# ═══════════════════════════════════════════════════════════

def select_top_rows(ai_long: pd.DataFrame, qid: str, top_k: int = 3, exclude_net: bool = True):
    """Select top K answer options for a given question ID."""
    dfq = ai_long[ai_long["question_id"] == qid].copy()
    if dfq.empty:
        return dfq

    if exclude_net and "is_net" in dfq.columns:
        non_net = dfq[dfq["is_net"] == False]
        if not non_net.empty:
            dfq = non_net

    if "rank_pct_desc" in dfq.columns:
        dfq = dfq.sort_values(["rank_pct_desc", "pct"], ascending=[True, False])
    else:
        dfq = dfq.sort_values(["pct"], ascending=False)

    return dfq.head(top_k)


def select_top_rows_multi(ai_long: pd.DataFrame, question_ids: list, top_k: int = 3, exclude_net: bool = True):
    """
    Select top K answer options across multiple questions.
    For multi-question slides, returns data grouped by question.
    """
    all_rows = []
    for qid in question_ids:
        rows = select_top_rows(ai_long, qid, top_k=top_k, exclude_net=exclude_net)
        if not rows.empty:
            all_rows.append(rows)

    if not all_rows:
        return pd.DataFrame()

    return pd.concat(all_rows, ignore_index=True)


def format_values(rows: pd.DataFrame, pct_decimals: int = 0):
    """Format rows as 'Option – XX%; Option – XX%; Option – XX%.'"""
    parts = []
    for _, r in rows.iterrows():
        opt = str(r["answer_option"]).strip()
        pct = float(r["pct"])
        if pct_decimals == 0:
            pct_str = f"{round(pct):.0f}%"
        else:
            pct_str = f"{pct:.{pct_decimals}f}%"
        parts.append(f"{opt} – {pct_str}")
    return "; ".join(parts) + "."


def format_values_grouped(ai_long: pd.DataFrame, question_ids: list, top_k: int = 3,
                           exclude_net: bool = True, pct_decimals: int = 0):
    """
    Format top K values for multiple questions, grouped by question.
    Returns a combined string with question labels.
    """
    parts = []
    for qid in question_ids:
        rows = select_top_rows(ai_long, qid, top_k=top_k, exclude_net=exclude_net)
        if rows.empty:
            continue
        row = rows.iloc[0]
        q_text = row.get("question_text", qid) if "question_text" in row.index else qid
        q_text = str(q_text).strip() if pd.notna(q_text) and str(q_text).strip() else qid
        # Shorten question text for display
        short_q = q_text[:80] + "..." if len(q_text) > 80 else q_text
        vals = format_values(rows, pct_decimals=pct_decimals)
        parts.append(f"{qid}: {vals}")

    return "\n".join(parts) if parts else ""


# ═══════════════════════════════════════════════════════════
#  Slide text extraction
# ═══════════════════════════════════════════════════════════

def get_slide_text(slide) -> str:
    """Extract all text from a slide."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if t:
                parts.append(t)
    return "\n".join(parts)


def get_question_text_from_slide(slide) -> str:
    """Get the question text from a multi-question/table slide (the shape with 'Questions X-Y:' or 'Question N:')."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if re.search(r"\bQuestions?\s+\d+(\s*[-–—]\s*\d+)?\s*:", text, re.IGNORECASE):
            return text
    return ""


def _find_restatement_shape(slide) -> str | None:
    """Find the shape with the chart callout/restatement on a question slide. Returns its text or None."""
    question_pattern = re.compile(r"\bQuestion[s]?\s+\d+", re.IGNORECASE)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or len(text) < 15:
            continue
        if question_pattern.search(text) and "Questions" not in text[:30]:
            continue  # Skip question shape
        if "%" in text and len(text) < 500:  # Summary-like: has %, not too long
            return text
    return None


def extract_chart_callouts_from_deck(prs, ai_long: pd.DataFrame) -> list[dict]:
    """
    Extract one-line chart callouts (restatements) from question slides.
    Returns list of {section, qid, question_text, summary}.
    """
    results = []
    for i, slide in enumerate(prs.slides):
        if is_section_divider(slide):
            continue
        slide_text = get_slide_text(slide)
        qspec = parse_question_spec(slide_text)
        if not qspec:
            continue
        qids = get_question_ids(qspec)
        if qspec[0] == "range" or len(qids) > 1:
            continue  # Multi-question: no single callout on slide
        summary = _find_restatement_shape(slide)
        if not summary:
            continue
        section = get_section_name_for_slide(prs, i)
        qid = qids[0]
        qt = ""
        rows = ai_long[ai_long["question_id"] == qid]
        if not rows.empty and "question_text" in rows.columns:
            qt = str(rows.iloc[0].get("question_text", "") or "").strip()
        results.append({"section": section, "qid": qid, "question_text": qt, "summary": summary})
    return results


def slide_has_placeholder(slide) -> bool:
    """Check if slide contains the placeholder text."""
    return PLACEHOLDER in get_slide_text(slide)


def _iter_all_shapes(container):
    """Yield all shapes including those inside groups (recursive)."""
    from pptx.shapes.group import GroupShape
    for shp in container.shapes:
        yield shp
        if isinstance(shp, GroupShape):
            yield from _iter_all_shapes(shp)


def _get_slide_number_shape(slide):
    """Return the slide number placeholder shape, or None. Searches top-level and inside groups."""
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
    for shp in _iter_all_shapes(slide):
        if not getattr(shp, "is_placeholder", False):
            continue
        if getattr(shp, "placeholder_format", None) is None:
            continue
        if getattr(shp.placeholder_format, "type", None) == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
            return shp
    return None


def _apply_arial_14_to_slide_number_shape(shape, force_white=False, force_not_bold=False):
    """Set Arial 14pt font on slide number placeholder via XML defRPr/endParaRPr.
    force_white: set text color to white (for blue/dark section divider slides).
    force_not_bold: ensure text is not bold (for blue slides where layout may have bold)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    if not getattr(shape, "has_text_frame", False):
        return
    txBody = shape.text_frame._txBody
    for p in txBody.findall(qn("a:p")):
        # Add or update defRPr for default font (Arial 14pt)
        pPr = p.find(qn("a:pPr"))
        if pPr is None:
            pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", "1400")  # 14pt
        defRPr.set("lang", "en-US")
        if force_not_bold:
            defRPr.set("b", "0")
        if force_white:
            solid_fill = defRPr.find(qn("a:solidFill"))
            if solid_fill is None:
                solid_fill = etree.SubElement(defRPr, qn("a:solidFill"))
            srgb = solid_fill.find(qn("a:srgbClr"))
            if srgb is None:
                srgb = etree.SubElement(solid_fill, qn("a:srgbClr"), val="FFFFFF")
            else:
                srgb.set("val", "FFFFFF")
        for tag in (qn("a:latin"), qn("a:ea"), qn("a:cs")):
            existing = defRPr.find(tag)
            if existing is not None:
                defRPr.remove(existing)
            etree.SubElement(defRPr, tag, typeface="Arial")
        # Also set endParaRPr if present (used by fields)
        endParaRPr = p.find(qn("a:endParaRPr"))
        if endParaRPr is not None:
            endParaRPr.set("sz", "1400")
            if force_not_bold:
                endParaRPr.set("b", "0")
            if force_white:
                solid_fill = endParaRPr.find(qn("a:solidFill"))
                if solid_fill is None:
                    solid_fill = etree.SubElement(endParaRPr, qn("a:solidFill"))
                srgb = solid_fill.find(qn("a:srgbClr"))
                if srgb is None:
                    etree.SubElement(solid_fill, qn("a:srgbClr"), val="FFFFFF")
                else:
                    srgb.set("val", "FFFFFF")
            for tag in (qn("a:latin"), qn("a:ea"), qn("a:cs")):
                existing = endParaRPr.find(tag)
                if existing is not None:
                    endParaRPr.remove(existing)
                etree.SubElement(endParaRPr, tag, typeface="Arial")


def _copy_slide_number_from_ref(ref_slide, dst_slide):
    """Copy slide number placeholder from ref_slide to dst_slide. Used for section dividers with wrong layout."""
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
    from pptx.oxml.ns import qn

    ref_shp = _get_slide_number_shape(ref_slide)
    if ref_shp is None:
        return False

    def _next_shape_id(sld):
        max_id = 0
        for shp in _iter_all_shapes(sld):
            try:
                sid = getattr(shp, "shape_id", None) or getattr(shp.element, "id", None)
                if sid is not None:
                    max_id = max(max_id, int(sid))
            except (TypeError, ValueError):
                pass
        return max_id + 1

    def _remap_embed_rids(el, src_part, dst_part):
        for e in el.iter():
            rid = e.get(qn("r:embed"))
            if rid and src_part.rels.get(rid):
                rel = src_part.rels[rid]
                new_rid = dst_part.rels.get_or_add(rel.reltype, rel.target_part)
                e.set(qn("r:embed"), new_rid)

    # Only use ref_slide—layout/master copy can corrupt (cross-part rIds)
    for src in [ref_slide]:
        if src is None:
            continue
        for shp in src.shapes:
            if not getattr(shp, "is_placeholder", False) or getattr(shp, "placeholder_format", None) is None:
                continue
            if getattr(shp.placeholder_format, "type", None) != PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
                continue
            new_el = copy.deepcopy(shp.element)
            nv_sp_pr = new_el.find(qn("p:nvSpPr"))
            if nv_sp_pr is not None:
                c_nv_pr = nv_sp_pr.find(qn("p:cNvPr"))
                if c_nv_pr is not None:
                    c_nv_pr.set("id", str(_next_shape_id(dst_slide)))
            if src is not ref_slide:
                _remap_embed_rids(new_el, src.part, dst_slide.part)
            sp_tree = dst_slide.shapes._spTree
            ext_lst = sp_tree.find(qn("p:extLst"))
            if ext_lst is not None:
                sp_tree.insert_element_before(new_el, qn("p:extLst"))
            else:
                sp_tree.append(new_el)
            return True
    return False


def _remove_slide_number_placeholder(slide):
    """Remove the slide number placeholder from a slide (including inside groups)."""
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
    from pptx.oxml.ns import qn

    for shp in list(_iter_all_shapes(slide)):
        if not getattr(shp, "is_placeholder", False) or getattr(shp, "placeholder_format", None) is None:
            continue
        if getattr(shp.placeholder_format, "type", None) != PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
            continue
        parent = shp.element.getparent()
        if parent is not None:
            parent.remove(shp.element)
            return True
    return False


def normalize_slide_numbers(prs, ref_slide=None):
    """
    Apply Arial 14pt to all slide number placeholders. Blue slides get white + not bold.
    Never touch position—preserves original layout. No remove/copy for blue slides (avoids corruption).
    """
    if ref_slide is None:
        for slide in prs.slides:
            if is_section_divider(slide):
                continue
            if _get_slide_number_shape(slide):
                ref_slide = slide
                break
    if ref_slide is None or _get_slide_number_shape(ref_slide) is None:
        return

    for slide in prs.slides:
        shp = _get_slide_number_shape(slide)
        # Only copy if slide has no placeholder (e.g. new transition slides)
        if shp is None:
            if _copy_slide_number_from_ref(ref_slide, slide):
                shp = _get_slide_number_shape(slide)
        if shp is None:
            continue
        # Apply font only—never touch position (preserves original placement, avoids corruption)
        is_divider = is_section_divider(slide)
        _apply_arial_14_to_slide_number_shape(
            shp, force_white=is_divider, force_not_bold=is_divider
        )


# ═══════════════════════════════════════════════════════════
#  Text replacement with formatting preservation
# ═══════════════════════════════════════════════════════════

def _copy_run_format(src_run, dst_run):
    """Copy font formatting from source run to destination run."""
    try:
        if src_run.font.name:
            dst_run.font.name = src_run.font.name
        if src_run.font.size:
            dst_run.font.size = src_run.font.size
        if src_run.font.bold is not None:
            dst_run.font.bold = src_run.font.bold
        if src_run.font.italic is not None:
            dst_run.font.italic = src_run.font.italic
        if src_run.font.color and src_run.font.color.rgb:
            dst_run.font.color.rgb = src_run.font.color.rgb
    except Exception:
        pass  # If we can't copy some attr, continue gracefully


def ensure_key_finding_style(prs: Presentation) -> dict:
    """
    Inspect the deck once to capture the font family, size, and color
    from the {Insert Finding Here} placeholder. This is the style we
    want for key-finding sentences and related bullets.
    """
    global KEY_FINDING_STYLE
    if KEY_FINDING_STYLE is not None:
        return KEY_FINDING_STYLE

    style: dict = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            if not tf.text or PLACEHOLDER not in tf.text:
                continue
            for para in tf.paragraphs:
                if para.runs:
                    r0 = para.runs[0]
                    f = r0.font
                    style = {
                        "name": f.name,
                        "size": f.size,
                        "bold": f.bold,
                        "italic": f.italic,
                        "rgb": getattr(getattr(f, "color", None), "rgb", None),
                    }
                    break
            if style:
                break
        if style:
            break

    KEY_FINDING_STYLE = style
    return KEY_FINDING_STYLE


def apply_style_to_run(run, style: dict, force_bold: bool | None = None) -> None:
    """Apply a captured style dict to a run. No-op if style is empty."""
    if not style:
        return
    f = run.font
    name = style.get("name")
    size = style.get("size")
    bold = style.get("bold")
    italic = style.get("italic")
    rgb = style.get("rgb")

    if name:
        f.name = name
    if size:
        f.size = size
    if force_bold is not None:
        f.bold = force_bold
    elif bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if rgb:
        try:
            f.color.rgb = rgb
        except Exception:
            # If color is theme-based or invalid, ignore
            pass


def _apply_bullet_formatting(paragraph, extra_marL_pt: float = 8):
    """
    Match transition slide bullet formatting: space between bullet and text (marL),
    and bullet color follows text (buClrTx).
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    p_el = paragraph._p
    pPr = p_el.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.SubElement(p_el, qn("a:pPr"))
    # Space between bullet and text (like transition slides)
    current = pPr.get("marL")
    current_emu = int(current) if current else 0
    extra_emu = int(Pt(extra_marL_pt))
    pPr.set("marL", str(current_emu + extra_emu))
    # Bullet color follows text (match transition slides)
    for tag in (qn("a:buClr"), qn("a:buClrTx")):
        for el in list(pPr.findall(tag)):
            pPr.remove(el)
    etree.SubElement(pPr, qn("a:buClrTx"))


def replace_placeholder_in_shape(shape, new_text: str) -> bool:
    """
    Replace {Insert Finding Here} placeholder with new_text,
    preserving the original text formatting as much as possible.
    Handles multi-line new_text by creating additional paragraphs.
    """
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame
    full_text = tf.text
    if PLACEHOLDER not in full_text:
        return False

    # Find the paragraph and run containing the placeholder
    for para_idx, para in enumerate(tf.paragraphs):
        para_text = para.text
        if PLACEHOLDER not in para_text:
            continue

        # Capture formatting from the first run (or default)
        font_name = None
        font_size = None
        font_bold = None
        font_italic = None
        font_color = None
        alignment = para.alignment

        if para.runs:
            ref_run = para.runs[0]
            font_name = ref_run.font.name
            font_size = ref_run.font.size
            font_bold = ref_run.font.bold
            font_italic = ref_run.font.italic
            try:
                font_color = ref_run.font.color.rgb
            except Exception:
                pass

        # Handle multi-line replacement
        lines = new_text.split("\n")

        # Clear ALL runs in this paragraph
        for run in list(para.runs):
            r_elem = run._r
            r_elem.getparent().remove(r_elem)

        # Minimal gap after line below heading (avoid pushing summary into chart)
        if para_idx == 0:
            # First para: use margin_top (space_before ignored by PowerPoint)
            try:
                current_pt = tf.margin_top.pt if tf.margin_top else 0
            except (AttributeError, TypeError):
                current_pt = 0
            tf.margin_top = Pt(current_pt + 6)
        else:
            para.space_before = Pt(6)

        # Set first line in the existing paragraph (leading space for bullet-text gap)
        para.level = 0
        run = para.add_run()
        run.text = " " + lines[0]
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        if font_bold is not None:
            run.font.bold = font_bold
        if font_italic is not None:
            run.font.italic = font_italic
        if font_color:
            run.font.color.rgb = font_color

        # Add additional paragraphs for remaining lines
        from pptx.oxml.ns import qn
        for line in lines[1:]:
            new_p = copy.deepcopy(para._p)
            # Clear runs in the copied paragraph
            for r in new_p.findall(qn('a:r')):
                new_p.remove(r)
            # Add a run with the text
            new_r = copy.deepcopy(run._r)
            new_r.text = " " + line
            new_p.append(new_r)
            # Insert after current paragraph
            para._p.addnext(new_p)

        _apply_bullet_formatting(para)
        return True

    return False


def set_shape_text_to_single_paragraph(shape, text: str, style: dict | None = None) -> bool:
    """
    Set the shape's text frame to a single paragraph with the given text.
    Removes all other paragraphs (avoids long stats overlapping tables).
    If a style dict is provided, apply it to the run (e.g. key-finding style).
    """
    if not shape.has_text_frame:
        return False

    from pptx.oxml.ns import qn

    tf = shape.text_frame
    txBody = tf._txBody

    # Keep only the first paragraph in the text frame
    paragraphs = list(txBody.findall(qn("a:p")))
    for p in paragraphs[1:]:
        txBody.remove(p)

    first_p = txBody.find(qn("a:p"))
    if first_p is None:
        return False

    # Remove all existing runs in the first paragraph
    for r in list(first_p.findall(qn("a:r"))):
        first_p.remove(r)

    p0 = tf.paragraphs[0]
    p0.level = 0  # Bullet level (like transition slides)
    run = p0.add_run()
    # Space between bullet and text (template provides bullet; " " matches transition spacing)
    run.text = " " + text

    if style:
        apply_style_to_run(run, style)

    # Minimal gap after line below heading (avoid pushing summary into chart)
    try:
        current_pt = tf.margin_top.pt if tf.margin_top else 0
    except (AttributeError, TypeError):
        current_pt = 0
    tf.margin_top = Pt(current_pt + 6)

    _apply_bullet_formatting(p0)
    return True


# ═══════════════════════════════════════════════════════════
#  Section detection
# ═══════════════════════════════════════════════════════════

def is_section_divider(slide) -> str | None:
    """
    Check if a slide is a section divider.
    Returns the section name if it is, None otherwise.

    Section dividers are identified by:
    1. Slides whose title matches a known SECTION_NAMES entry
    2. OR slides that look like dividers: single title, no chart/table, not a question slide
    """
    texts = []
    has_chart = any(getattr(s, "has_chart", False) for s in slide.shapes)
    has_table = any(getattr(s, "has_table", False) for s in slide.shapes)

    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and not t.isdigit():  # Ignore slide numbers
                texts.append(t)

    if not texts:
        return None

    # 1. Check if any text matches a known section name
    for text in texts:
        clean = text.strip()
        for section in SECTION_NAMES:
            if clean.lower() == section.lower():
                return section

    # 2. Dynamic detection: single title only, no chart/table, not a question slide
    if has_chart or has_table:
        return None
    if len(texts) != 1:  # Divider has exactly one text block (section name only)
        return None
    primary = texts[0].strip()
    if not primary or len(primary) > 80:  # Too long for a section title
        return None
    if parse_question_spec(primary) is not None:  # "Question 5:" etc.
        return None
    if re.search(r"Q\d+", primary, re.IGNORECASE):  # Contains Q7, Q8, etc.
        return None
    return primary


def get_section_questions(prs, section_slide_idx: int, ai_long: pd.DataFrame) -> dict:
    """
    Get all questions belonging to a section.
    A section spans from the section divider slide to the next section divider (or end).

    Returns dict with:
        - question_ids: list of question IDs
        - question_data: DataFrame of all data for these questions
        - question_texts: dict mapping qid -> question_text
    """
    question_ids = []
    slides = list(prs.slides)

    # Scan forward from section divider until next divider or end
    for i in range(section_slide_idx + 1, len(slides)):
        slide = slides[i]

        # Stop if we hit another section divider
        if is_section_divider(slide):
            break

        slide_text = get_slide_text(slide)
        qspec = parse_question_spec(slide_text)
        if qspec:
            qids = get_question_ids(qspec)
            question_ids.extend(qids)

    # Remove duplicates while preserving order
    seen = set()
    unique_qids = []
    for qid in question_ids:
        if qid not in seen:
            seen.add(qid)
            unique_qids.append(qid)

    # Get data
    question_data = ai_long[ai_long["question_id"].isin(unique_qids)].copy()

    # Get question texts (fallback to qid if column missing or empty)
    question_texts = {}
    for qid in unique_qids:
        rows = ai_long[ai_long["question_id"] == qid]
        if not rows.empty:
            row = rows.iloc[0]
            qt = row.get("question_text", qid) if "question_text" in row.index else qid
            question_texts[qid] = str(qt).strip() if pd.notna(qt) and str(qt).strip() else qid

    return {
        "question_ids": unique_qids,
        "question_data": question_data,
        "question_texts": question_texts,
    }


def get_section_name_for_slide(prs, slide_idx: int) -> str:
    """
    Get the section name for a slide at the given index.
    Looks backwards for the nearest section divider.
    """
    slides = list(prs.slides)
    for i in range(slide_idx, -1, -1):
        sec = is_section_divider(slides[i])
        if sec:
            return sec
    return "Key Findings"


# ═══════════════════════════════════════════════════════════
#  LLM helpers
# ═══════════════════════════════════════════════════════════

def call_llm(system_prompt: str, user_prompt: str, model: str = "gpt-4.1-mini", temperature: float = 0.2) -> str:
    """Call OpenAI-compatible LLM. Returns response text."""
    import os
    from openai import OpenAI

    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        raise ValueError("OPENAI_API_KEY environment variable not set")

    client = OpenAI(api_key=api_key)
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=temperature,
    )
    return resp.choices[0].message.content.strip()


def generate_restatement(bullets: str, question_context: str | None = None) -> str:
    """Generate a one-line chart callout that incorporates the question and editorializes in the client's favor."""
    system = (
        "You are writing a one-line KEY FINDING callout for a professional polling deck. "
        "CRITICAL: Echo the question in your sentence—weave in what was asked so the finding reads in context. "
        "Write in past tense (e.g., '32% said' not '32% say'). "
        "Use a slightly less formal, conversational tone. "
        "Avoid generic phrases like 'the survey showed', 'overall ratings', 'respondents indicated' without the question context. "
        "Editorialize subtly in the client's favor where appropriate."
    )

    question_block = ""
    if question_context and question_context.strip():
        question_block = f"""
The survey question (MUST echo this in your summary):
\"\"\"{question_context.strip()}\"\"\"

GOOD: "In a typical week, 32% said they rode the New York City subway or buses 3-4 days per week, and 26% rode 5 or more days per week."
BAD: "The survey showed a nearly even split in overall ratings, with about half giving fair or poor assessments."
BAD: "The most common response was 3-4 days per week."
"""

    user = f"""
Write EXACTLY ONE sentence summarizing these survey results.
{question_block}
Rules:
- ECHO the question: start with the question context (e.g., "In a typical week, 32% said they rode...", "When asked about favorability, 60% said...")
- NO generic language: avoid "the survey showed", "overall ratings", "respondents indicated", "the most common response is"
- Write in past tense. Use a slightly less formal tone.
- Lead with the top results and percentages. Editorialize subtly in the client's favor.
- Maximum 35 words. Do NOT invent numbers.

Survey results:
{bullets}

Output ONLY the sentence.
"""

    out = call_llm(system, user)
    return _strip_question_ids(out)


def generate_questions_asked_content(section_name: str, question_texts: dict) -> str:
    """Generate 'Questions Asked' transition slide content."""
    q_list = "\n".join([f"- {qid}: {text}" for qid, text in question_texts.items()])

    system = (
        "You are a political survey analyst writing for senior decision-makers. "
        "Describe what we were trying to learn and what questions we asked. "
        "Focus on the TOPICS and CONCEPTS, not methodology. "
        "Write in plain language—avoid technical jargon like 'scaled responses', 'frequency-based items', 'measurement refinement'. "
        "Write in past tense."
    )
    user = f"""You are preparing a transition slide for the "{section_name}" section of a survey deck.

Section name: "{section_name}"

Questions in this section (IDs and abbreviated text):
{q_list}

Write 5–6 bullet points that describe:
1. **What we were trying to learn** – the topics, concepts, or issues we explored (e.g., favorability, ridership patterns, voter concerns).
2. **What questions we asked** – summarize the substance of the questions in plain language (e.g., "We asked about transit use and how often people ride the subway").
3. **Why it matters** – how the answers help inform strategy or decisions.
4. **What to expect** – how these findings set up what appears in the next slides.

Guidelines:
- Write in past tense. Use plain, conversational language.
- NO technical jargon: avoid "scaled responses", "frequency-based items", "measurement refinement", "dimensions/topics being probed".
- NO numeric values. NEVER use question numbers (Q7, Q10, etc.); use the question topic/concept instead.
- Write as clean bullet points (start each line with "• ").
- Keep total length under 900 characters."""

    out = call_llm(system, user)
    return _strip_question_ids(out)


def generate_survey_responses_content(section_name: str, question_texts: dict,
                                       question_data: pd.DataFrame) -> str:
    """Generate 'Survey Responses' transition slide content."""
    # Build data summary for each question (include full question context for LLM to reference)
    data_parts = []
    for qid, text in question_texts.items():
        qdata = question_data[question_data["question_id"] == qid]
        if qdata.empty:
            continue
        top = qdata.sort_values("pct", ascending=False).head(3)
        vals = "; ".join([f"{r['answer_option'].strip()} – {r['pct']:.0f}%" for _, r in top.iterrows()])
        data_parts.append(f"[{qid}] Question: {text[:250]}\n  Top results: {vals}")

    data_summary = "\n".join(data_parts)

    system = (
        "You are a political survey analyst preparing executive briefing slides. "
        "Summarize survey RESULTS by echoing what was asked in each finding. "
        "Every bullet must weave in the question context—e.g., 'When asked about X, 60% said Y' NOT 'The survey showed a nearly even split'. "
        "Avoid generic phrases like 'overall ratings', 'respondents gave fair or poor assessments' without saying what was asked. "
        "Write in past tense. Stay strictly faithful to the numbers provided."
    )
    user = f"""You are preparing a transition slide for the "{section_name}" section of a survey deck.

Section name: "{section_name}"

Data for this section (per question, with top answer options and percentages).
The [Q7], [Q10] etc. IDs are for your reference only—do NOT include them in your output.
{data_summary}

Write 5–6 bullet points. CRITICAL: Each bullet must ECHO what was asked.

GOOD: "When asked to rate their transit experience, 65% said good or excellent and 35% said fair or poor."
BAD: "The survey showed a nearly even split in overall ratings, with about half giving fair or poor assessments."

GOOD: "When asked about favorability of Tom Malinowski, 60% said very or somewhat favorable."
BAD: "Respondents gave mixed assessments across the board."

Guidelines:
- EVERY bullet must weave in the question context. Start with "When asked about...", "Regarding...", or similar so readers know what was asked.
- NO generic language: avoid "the survey showed", "overall ratings", "respondents indicated" without the question context.
- NO question numbers (Q7, Q10). Use the topic/concept from the question text.
- Use key percentages from the data where helpful (e.g., “around 6 in 10”, or explicit % when clear).
- Do NOT invent any new numbers or options beyond what appears in the data summary above.
- Write as clean bullet points (start each line with "• ").
- Keep total length under 900 characters."""

    out = call_llm(system, user)
    return _strip_question_ids(out)


def generate_multi_question_summary_content(section_name: str, question_texts: dict,
                                            question_data: pd.DataFrame) -> str:
    """
    Generate summary content for a multi-question slide: what we asked and what we concluded.
    Same structure as key findings / transition slides: bullets with question context + results.
    """
    q_list = "\n".join([f"- {qid}: {text}" for qid, text in question_texts.items()])
    data_parts = []
    for qid, text in question_texts.items():
        qdata = question_data[question_data["question_id"] == qid]
        if qdata.empty:
            continue
        top = qdata.sort_values("pct", ascending=False).head(3)
        vals = "; ".join([f"{r['answer_option'].strip()} – {r['pct']:.0f}%" for _, r in top.iterrows()])
        data_parts.append(f"[{qid}] Question: {text[:250]}\n  Top results: {vals}")
    data_summary = "\n".join(data_parts)

    system = (
        "You are a political survey analyst preparing executive briefing slides. "
        "Write a summary slide for a multi-question slide: what we asked and what we concluded. "
        "EVERY bullet must echo what was asked—weave in the question context (e.g., 'When asked about X, 60% said Y'). "
        "Avoid generic language like 'the survey showed' or 'respondents indicated' without the question context. "
        "Write in past tense. Use a slightly less formal tone. Stay strictly faithful to the numbers."
    )
    user = f"""You are preparing a summary slide for a multi-question slide in the "{section_name}" section.

Section: "{section_name}"

Questions on this slide (IDs for reference only—do NOT include them in output):
{q_list}

Data (top results per question):
{data_summary}

Write 4–6 bullet points that:
1. Summarize what we asked (weave in question context, e.g. "When asked about X, 60% said very favorable")
2. Summarize what we concluded (key findings with percentages)
3. Use past tense throughout (e.g., 'we asked', 'respondents said', 'the survey showed'). Use a slightly less formal tone. Editorialize subtly in the client's favor where the data supports it.
4. NEVER use question numbers (Q7, Q10) in output—use the question topic/concept instead.
5. Write as clean bullet points (start each line with "• ").
6. Keep total length under 900 characters."""

    out = call_llm(system, user)
    return _strip_question_ids(out)


# ═══════════════════════════════════════════════════════════
#  Executive Summary
# ═══════════════════════════════════════════════════════════

DEFAULT_EXEC_SUMMARY_PROMPT = """You are writing a 10-15-slide executive summary for a PowerPoint about a public opinion survey. The audience is senior decision-makers.

The executive summary needs to provide a high-level summary of the key findings from the survey research as well as providing actionable insights in terms of what the data means and what actions should be taken based on the findings of the survey. The executive summary should end with 1 page of strategic recommendations that are based on the research finding and are tangible next steps that should be taken.

Use this structure:
1) Goals/objectives (1 slide) – what did we seek to accomplish with the surveying?
2) Narrative in One Slide (1 slide) – what is the high-level elevator pitch of what should be the strategic pathway forward?
3) Main conclusions by theme (7-12 slides) – require that every theme slide ends in an implication. That's what turns polling into strategy.
4) Key Themes to drive (1 slide) – no need to present more data – what are the key themes that emerge from the survey that are on people's minds?
5) Strategy & tactics (2 slides)

Rules:
- Each slide must have: Slide title (conclusion headline) + 3-6 bullets.
- Keep slide titles to 8-9 words maximum so they fit on one line.
- Every theme slide must include: Insight → proof (2-4 stats) → implication.
- Use plain language. Avoid methodological jargon unless necessary.
- Keep it persuasive: tell a story arc that leads to clear actions.
- Write in the voice of a strategic pollster: confident, concise, practical. No hype.
- Write in past tense.
- NEVER use question numbers (Q7, Q10, etc.) – use the question topic/concept instead.
- Replace long dashes (—) with regular hyphens (-) or commas."""


def generate_executive_summary_slides(
    goals_text: str,
    sections_data: list[dict],
    ai_long: pd.DataFrame,
    summaries_text: str | None = None,
) -> list[dict]:
    """
    Generate 12-15 executive summary slides.
    Uses summaries_text (chart callouts, transition slides, multi-q summaries) when provided.
    Returns list of {"title": str, "bullets": list[str]}.
    """
    # Primary input: pre-generated summaries from the deck (chart callouts, transition slides, multi-q)
    if summaries_text and summaries_text.strip():
        data_block = f"""PRE-GENERATED SUMMARIES FROM THIS DECK (use these as your primary source; synthesize and elevate):
{summaries_text}
"""
    else:
        data_block = ""

    # Fallback: raw survey data by section (when summaries unavailable)
    data_parts = []
    for sec in sections_data:
        name = sec.get("name", "")
        q_data = sec.get("questions", {})
        q_texts = q_data.get("question_texts", {})
        q_ids = q_data.get("question_ids", [])
        question_data = q_data.get("question_data", pd.DataFrame())
        if question_data.empty and "question_id" in ai_long.columns:
            question_data = ai_long[ai_long["question_id"].isin(q_ids)].copy()
        lines = [f"Section: {name}"]
        for qid, text in q_texts.items():
            qrows = question_data[question_data["question_id"] == qid] if not question_data.empty else pd.DataFrame()
            if qrows.empty:
                qrows = ai_long[ai_long["question_id"] == qid]
            top = qrows.sort_values("pct", ascending=False).head(3) if not qrows.empty else pd.DataFrame()
            vals = "; ".join([f"{r['answer_option'].strip()} – {r['pct']:.0f}%" for _, r in top.iterrows()]) if not top.empty else "no data"
            lines.append(f"  - Question: {text[:120]}...")
            lines.append(f"    Top results: {vals}")
        data_parts.append("\n".join(lines))
    survey_data_summary = "\n\n".join(data_parts)

    system = (
        "You are a strategic pollster writing an executive summary for senior decision-makers. "
        "Output MUST follow the exact format below. Each slide block starts with '=== SLIDE N: Type ===' "
        "then 'Title: [conclusion headline]' on the next line, then bullet points each starting with '• '."
    )
    user = f"""{DEFAULT_EXEC_SUMMARY_PROMPT}

SPECIFIC GOALS OF THIS SURVEY:
{goals_text}

{data_block}
SURVEY DATA (by section, for reference):
{survey_data_summary}

OUTPUT FORMAT – use EXACTLY this structure. Output 12-15 slides.

=== SLIDE 1: Goals ===
Title: [Your conclusion headline for goals - max 8-9 words]
• [bullet 1]
• [bullet 2]
• [bullet 3]

=== SLIDE 2: Narrative ===
Title: [Your elevator pitch headline - max 8-9 words]
• [bullet 1]
• [bullet 2]
...

=== SLIDE 3: Theme 1 ===
Title: [Conclusion headline for first theme - max 8-9 words]
• [Insight]
• [Proof with 2-4 stats]
• [Implication – what to do next]

... (continue for 7-12 theme slides, one per major section/finding)

=== SLIDE N: Key Themes ===
Title: [Headline]
• [theme 1]
• [theme 2]
...

=== SLIDE N+1: Strategy ===
Title: [Headline]
• [tactic 1]
• [tactic 2]
...

=== SLIDE N+2: Tactics ===
Title: [Headline]
• [tactic 1]
• [tactic 2]
...

Generate the full executive summary now. Synthesize and elevate the pre-generated summaries when provided; otherwise use the survey data. No question numbers."""

    out = call_llm(system, user, temperature=0.3)
    out = _strip_question_ids(out)
    # Replace em dashes with regular hyphen
    out = out.replace("—", "-")

    return _parse_executive_summary_output(out)


def _parse_executive_summary_output(text: str) -> list[dict]:
    """Parse LLM output into list of {title, bullets}. Use structure titles so readers understand purpose."""
    raw_slides = []
    blocks = re.split(r"=== SLIDE \d+[^=]*===", text, flags=re.IGNORECASE)
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        bullets = []
        for line in block.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.lower().startswith("title:"):
                pass
            elif line.startswith("•") or line.startswith("-") or line.startswith("*"):
                bullets.append(line.lstrip("•-* ").strip())
        if bullets:
            raw_slides.append({"bullets": bullets})
    raw_slides = raw_slides[:15]
    # Assign structure titles: Goals → Narrative → Theme 1..N → Key Themes → Strategy → Tactics
    n = len(raw_slides)
    if n == 0:
        return []
    num_themes = max(0, n - 5)  # 5 fixed: Goals, Narrative, Key Themes, Strategy, Tactics
    titles = ["Goals/objectives", "Narrative in One Slide"]
    titles.extend(f"Theme {i + 1}" for i in range(num_themes))
    titles.extend(["Key Themes to drive", "Strategy", "Tactics"])
    titles = titles[:n]  # Trim if fewer slides
    while len(titles) < n:
        titles.append(f"Theme {len(titles) - 1}")
    return [{"title": titles[i], "bullets": raw_slides[i]["bullets"]} for i in range(n)]
