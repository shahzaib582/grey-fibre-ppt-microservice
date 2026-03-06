"""
PASS 1 — Insert Top 3 Numeric Findings
=======================================
For each slide containing {Insert Finding Here}:
  - Detect Question number (single or range)
  - Pull Top 3 answer options (excluding NETs if flagged)
  - Insert formatted values: Option – XX%; Option – XX%; Option – XX%.

Usage:
  python -m survey_pipeline.pass1_insert_numbers --data DATA.xlsx --pptx TEMPLATE.pptx --out output_pass1.pptx
"""

import argparse
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from .utils import (
    PLACEHOLDER,
    parse_question_spec,
    get_question_ids,
    select_top_rows,
    format_values,
    format_values_grouped,
    get_slide_text,
    replace_placeholder_in_shape,
)
from .data_loader import load_ai_long


def _update_chart_for_single_question(slide, rows):
    """
    Update a simple category chart on the slide using the given rows.

    Assumes a single-series bar/column chart where categories correspond
    to answer options and values to percentages.
    """
    categories = [str(r["answer_option"]).strip() for _, r in rows.iterrows()]
    values = [float(r["pct"]) for _, r in rows.iterrows()]

    if not categories or not values:
        return False

    for shape in slide.shapes:
        # Older python-pptx exposes has_chart on GraphicFrame shapes
        if not hasattr(shape, "has_chart") or not shape.has_chart:
            continue

        chart = shape.chart
        # Limit to standard column/bar charts to avoid breaking complex visuals
        if chart.chart_type not in (
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            XL_CHART_TYPE.COLUMN_STACKED,
            XL_CHART_TYPE.BAR_CLUSTERED,
            XL_CHART_TYPE.BAR_STACKED,
        ):
            continue

        chart_data = CategoryChartData()
        chart_data.categories = categories

        if chart.series:
            series_name = chart.series[0].name
        else:
            series_name = "Series 1"

        chart_data.add_series(series_name, values)
        chart.replace_data(chart_data)
        return True

    return False


def _update_chart_for_multi_questions(slide, ai_long, qids, exclude_net=True):
    """
    Update a multi-series category chart for a range of questions.

    Assumes:
      - Categories correspond to answer options (e.g. Very / Somewhat ...)
      - Each series corresponds to one question ID (e.g. Q6, Q7, ...)
    """
    # Build per-question distributions (answer_option -> pct)
    series_list = []
    categories = None

    for qid in qids:
        dfq = ai_long[ai_long["question_id"] == qid].copy()
        if dfq.empty:
            continue

        if exclude_net and "is_net" in dfq.columns:
            non_net = dfq[dfq["is_net"] == False]
            if not non_net.empty:
                dfq = non_net

        if "rank_pct_desc" in dfq.columns:
            dfq = dfq.sort_values(["rank_pct_desc", "pct"], ascending=[True, False])
        else:
            dfq = dfq.sort_values(["pct"], ascending=False)

        # Use the first question as the category order
        if categories is None:
            categories = [str(a).strip() for a in dfq["answer_option"].tolist()]

        # Align this question's values to the shared category list
        cat_to_pct = {
            str(row["answer_option"]).strip(): float(row["pct"])
            for _, row in dfq.iterrows()
        }
        values = [cat_to_pct.get(cat, 0.0) for cat in categories]
        series_list.append((qid, values))

    if not series_list or not categories:
        return False

    for shape in slide.shapes:
        if not hasattr(shape, "has_chart") or not shape.has_chart:
            continue

        chart = shape.chart
        if chart.chart_type not in (
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            XL_CHART_TYPE.COLUMN_STACKED,
            XL_CHART_TYPE.BAR_CLUSTERED,
            XL_CHART_TYPE.BAR_STACKED,
        ):
            continue

        chart_data = CategoryChartData()
        chart_data.categories = categories

        for series_name, values in series_list:
            chart_data.add_series(series_name, values)

        chart.replace_data(chart_data)
        return True

    return False


def process_slide(slide, ai_long, top_k=3, exclude_net=True, pct_decimals=0):
    """Process a single slide: detect question, format values, replace placeholder."""
    slide_text = get_slide_text(slide)

    if PLACEHOLDER not in slide_text:
        return False

    qspec = parse_question_spec(slide_text)
    if qspec is None:
        print(f"  [WARN] Slide has placeholder but no question detected")
        return False

    qids = get_question_ids(qspec)
    if not qids:
        return False

    chart_updated = False
    fallback = "No data available for this question."

    if qspec[0] != "single":
        # For now we leave multi-question (range) slides completely untouched.
        print(f"  [SKIP] Multi-question slide {qids} — no numeric insert or chart update")
        return False

    # Single question — top K from that question
    rows = select_top_rows(ai_long, qids[0], top_k=top_k, exclude_net=exclude_net)
    if rows.empty:
        print(f"  [WARN] No data found for {qids[0]} — using fallback text")
        values_text = fallback
    else:
        values_text = format_values(rows, pct_decimals=pct_decimals)
        if _update_chart_for_single_question(slide, rows):
            chart_updated = True

    # Replace placeholder in all shapes
    replaced = False
    for shape in slide.shapes:
        if replace_placeholder_in_shape(shape, values_text):
            replaced = True

    if replaced:
        if chart_updated:
            print(f"  [OK] {qspec[0].upper()} {qids} → inserted values + chart")
        elif values_text == fallback:
            print(f"  [OK] {qspec[0].upper()} {qids} → fallback (no data)")
        else:
            print(f"  [OK] {qspec[0].upper()} {qids} → inserted values")

    return replaced


def main():
    ap = argparse.ArgumentParser(description="Pass 1: Insert Top 3 numeric findings")
    ap.add_argument(
        "--data",
        required=True,
        help="Path to survey Excel file (either AI-ready with 'ai_long' sheet or raw 250870-style 'ExcelData')",
    )
    ap.add_argument("--pptx", required=True, help="Path to PowerPoint template")
    ap.add_argument("--out", required=True, help="Output file path")
    ap.add_argument("--exclude-net", action="store_true", default=True, help="Exclude NET aggregates")
    ap.add_argument("--top-k", type=int, default=3, help="Number of top answers per question")
    ap.add_argument("--pct-decimals", type=int, default=0, help="Decimal places for percentages")
    args = ap.parse_args()

    print("=" * 60)
    print("PASS 1 — Insert Top 3 Numeric Findings")
    print("=" * 60)
    print(f"Data:     {args.data}")
    print(f"Template: {args.pptx}")
    print(f"Output:   {args.out}")
    print()

    ai_long = load_ai_long(args.data)
    print(f"Loaded {len(ai_long)} rows of survey data")
    try:
        qids_sorted = sorted(ai_long["question_id"].unique(), key=lambda x: int(str(x)[1:]))
        print(f"Questions: {qids_sorted}")
    except Exception:
        pass
    print()

    prs = Presentation(args.pptx)
    updated_count = 0
    total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        if process_slide(slide, ai_long, top_k=args.top_k,
                         exclude_net=args.exclude_net, pct_decimals=args.pct_decimals):
            updated_count += 1

    prs.save(args.out)
    print()
    print(f"PASS 1 complete.")
    print(f"  Updated: {updated_count} slide(s)")
    print(f"  Total:   {total_slides} slide(s)")
    print(f"  Output:  {args.out}")


if __name__ == "__main__":
    main()
