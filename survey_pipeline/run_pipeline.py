"""
Pipeline Orchestrator — Run All 3 Passes
==========================================
Runs the complete 3-step survey slide automation pipeline:

  Step 1: Insert Top 3 numeric findings (pass1_insert_numbers.py)
  Step 2: Add AI restatement sentences (pass2_add_restatement.py)
  Step 3: Generate transition slides (pass3_transition_slides.py)

Usage:
  python -m survey_pipeline.run_pipeline --data DATA.xlsx --pptx TEMPLATE.pptx --out FINAL.pptx

  # Run specific passes only:
  python -m survey_pipeline.run_pipeline --data DATA.xlsx --pptx TEMPLATE.pptx --out FINAL.pptx --passes 1
  python -m survey_pipeline.run_pipeline --data DATA.xlsx --pptx TEMPLATE.pptx --out FINAL.pptx --passes 1,2
"""

import argparse
import os
import sys
import tempfile
import shutil
from datetime import datetime


def run_pass1(data_path, pptx_path, out_path, top_k=3, pct_decimals=0):
    """Run Pass 1: Insert top 3 numeric findings."""
    from survey_pipeline.pass1_insert_numbers import main as pass1_main
    sys.argv = [
        "pass1_insert_numbers.py",
        "--data", data_path,
        "--pptx", pptx_path,
        "--out", out_path,
        "--exclude-net",
        "--top-k", str(top_k),
        "--pct-decimals", str(pct_decimals),
    ]
    pass1_main()


def run_pass2(data_path, pptx_path, out_path, top_k=3):
    """Run Pass 2: Add AI restatement sentences."""
    from survey_pipeline.pass2_add_restatement import main as pass2_main
    sys.argv = [
        "pass2_add_restatement.py",
        "--data", data_path,
        "--pptx", pptx_path,
        "--out", out_path,
        "--exclude-net",
        "--top-k", str(top_k),
    ]
    pass2_main()


def run_pass3(data_path, pptx_path, out_path):
    """Run Pass 3: Generate transition slides."""
    from survey_pipeline.pass3_transition_slides import main as pass3_main
    sys.argv = [
        "pass3_transition_slides.py",
        "--data", data_path,
        "--pptx", pptx_path,
        "--out", out_path,
    ]
    pass3_main()


def validate_output(pptx_path, original_slide_count, num_sections):
    """Validate the final output file."""
    from pptx import Presentation
    from survey_pipeline.utils import PLACEHOLDER, is_section_divider

    print("\n" + "=" * 60)
    print("VALIDATION")
    print("=" * 60)

    prs = Presentation(pptx_path)
    issues = []

    # Helper: get a simple title string for a slide (first non-empty text)
    def _get_title(slide):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    return text
        return ""

    # Helper: get all text for character-count checks
    def _get_all_text(slide):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t:
                    parts.append(t)
        return "\n".join(parts)

    # Check 1: No remaining placeholders
    placeholder_slides = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and PLACEHOLDER in shape.text_frame.text:
                placeholder_slides.append(i + 1)
    if placeholder_slides:
        issues.append(f"Remaining placeholders on slides: {placeholder_slides}")
    else:
        print("  [OK] All placeholders replaced")

    # Check 2: Slide count
    expected = original_slide_count + (2 * num_sections)
    actual = len(prs.slides)
    if actual == expected:
        print(f"  [OK] Slide count correct: {actual}")
    else:
        issues.append(f"Slide count: expected {expected}, got {actual}")
        print(f"  [!] Slide count: expected {expected}, got {actual}")

    # Check 3: Transition slides per section (only if Pass 3 ran)
    if num_sections > 0:
        # Find section divider indices in the FINAL deck
        section_indices = []
        for idx, slide in enumerate(prs.slides):
            sec_name = is_section_divider(slide)
            if sec_name:
                section_indices.append((idx, sec_name))

        if len(section_indices) != num_sections:
            issues.append(
                f"Section divider count mismatch in final deck: expected {num_sections}, "
                f"found {len(section_indices)}"
            )
        else:
            for j, (start_idx, section_name) in enumerate(section_indices):
                end_idx = section_indices[j + 1][0] if j + 1 < len(section_indices) else len(prs.slides)

                qa_title = f"{section_name}: Questions Asked"
                sr_title = f"{section_name}: Survey Responses"
                qa_slides = []
                sr_slides = []

                for k in range(start_idx + 1, end_idx):
                    title = _get_title(prs.slides[k])
                    if title == qa_title:
                        qa_slides.append(k + 1)  # human 1-based
                    elif title == sr_title:
                        sr_slides.append(k + 1)

                if len(qa_slides) != 1 or len(sr_slides) != 1:
                    issues.append(
                        f"Section '{section_name}' expects 1 'Questions Asked' and 1 'Survey Responses' slide; "
                        f"found QA={len(qa_slides)}, SR={len(sr_slides)} "
                        f"(QA slides: {qa_slides or 'none'}, SR slides: {sr_slides or 'none'})"
                    )
                else:
                    print(
                        f"  [OK] Section '{section_name}': "
                        f"transition slides on slides {qa_slides[0]} (Questions Asked) "
                        f"and {sr_slides[0]} (Survey Responses)"
                    )

                # Character-count check (<= 1000 chars) on transition slides only
                for slide_num in qa_slides + sr_slides:
                    text_len = len(_get_all_text(prs.slides[slide_num - 1]))
                    if text_len > 1000:
                        issues.append(
                            f"Transition slide {slide_num} for section '{section_name}' "
                            f"exceeds 1000 characters (len={text_len})"
                        )

    # Check 4: File opens without error
    try:
        _ = Presentation(pptx_path)
        print(f"  [OK] File opens without corruption")
    except Exception as e:
        issues.append(f"File corruption: {e}")

    if issues:
        print(f"\n  [!] {len(issues)} issue(s) found:")
        for issue in issues:
            print(f"    - {issue}")
    else:
        print(f"\n  [OK] All validations passed!")

    return len(issues) == 0


def main():
    ap = argparse.ArgumentParser(
        description="3-Step Survey Slide Automation Pipeline",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Full pipeline:
  python -m survey_pipeline.run_pipeline --data Wallenstein_Weighted_Interview_AI_Ready.xlsx --pptx Brendan_Gill_Slide_ExportPPT_1_23_2026.pptx --out Brendan_Gill_Enhanced_Final.pptx

  # Pass 1 only (no API key needed):
  python -m survey_pipeline.run_pipeline --data Wallenstein_Weighted_Interview_AI_Ready.xlsx --pptx Brendan_Gill_Slide_ExportPPT_1_23_2026.pptx --out output_pass1.pptx --passes 1

  # Passes 1 & 2:
  python -m survey_pipeline.run_pipeline --data Wallenstein_Weighted_Interview_AI_Ready.xlsx --pptx Brendan_Gill_Slide_ExportPPT_1_23_2026.pptx --out output_pass2.pptx --passes 1,2
        """,
    )
    ap.add_argument(
        "--data",
        required=True,
        help="Path to survey Excel (either AI-ready with 'ai_long' sheet or raw 250870-style 'ExcelData')",
    )
    ap.add_argument("--pptx", required=True, help="Path to PowerPoint template")
    ap.add_argument("--out", required=True, help="Final output file path")
    ap.add_argument("--passes", default="1,2,3", help="Comma-separated pass numbers to run (default: 1,2,3)")
    ap.add_argument("--top-k", type=int, default=3, help="Number of top answers per question")
    ap.add_argument("--pct-decimals", type=int, default=0, help="Decimal places for percentages")
    ap.add_argument("--keep-intermediates", action="store_true", help="Keep intermediate files")
    args = ap.parse_args()

    passes = [int(p.strip()) for p in args.passes.split(",")]

    print("+" + "=" * 58 + "+")
    print("|   3-STEP SURVEY SLIDE AUTOMATION PIPELINE                |")
    print("+" + "=" * 58 + "+")
    print(f"  Data:     {args.data}")
    print(f"  Template: {args.pptx}")
    print(f"  Output:   {args.out}")
    print(f"  Passes:   {passes}")
    print(f"  Started:  {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print()

    # Validate inputs
    if not os.path.exists(args.data):
        print(f"ERROR: Data file not found: {args.data}")
        sys.exit(1)
    if not os.path.exists(args.pptx):
        print(f"ERROR: Template file not found: {args.pptx}")
        sys.exit(1)

    # Get original slide count for validation
    from pptx import Presentation
    from survey_pipeline.utils import is_section_divider
    prs = Presentation(args.pptx)
    original_slide_count = len(prs.slides)
    num_sections = sum(1 for slide in prs.slides if is_section_divider(slide))
    del prs

    # Create temp directory for intermediate files
    base_dir = os.path.dirname(os.path.abspath(args.out))
    temp_dir = tempfile.mkdtemp(dir=base_dir, prefix="pipeline_temp_")

    try:
        current_pptx = args.pptx

        # ---- Pass 1 ----
        if 1 in passes:
            pass1_out = os.path.join(temp_dir, "output_pass1.pptx")
            print("\n" + "-" * 60)
            run_pass1(args.data, current_pptx, pass1_out,
                      top_k=args.top_k, pct_decimals=args.pct_decimals)
            current_pptx = pass1_out

            if args.keep_intermediates:
                shutil.copy2(pass1_out, os.path.join(base_dir, "output_pass1.pptx"))

        # ---- Pass 2 ----
        if 2 in passes:
            if "OPENAI_API_KEY" not in os.environ:
                print("\n[!] WARNING: OPENAI_API_KEY not set. Pass 2 requires an API key.")
                print("  Set it with: set OPENAI_API_KEY=your-key-here")
                print("  Skipping Pass 2.\n")
            else:
                pass2_out = os.path.join(temp_dir, "output_pass2.pptx")
                print("\n" + "-" * 60)
                run_pass2(args.data, current_pptx, pass2_out, top_k=args.top_k)
                current_pptx = pass2_out

                if args.keep_intermediates:
                    shutil.copy2(pass2_out, os.path.join(base_dir, "output_pass2.pptx"))

        # ---- Pass 3 ----
        if 3 in passes:
            if "OPENAI_API_KEY" not in os.environ:
                print("\n[!] WARNING: OPENAI_API_KEY not set. Pass 3 requires an API key.")
                print("  Set it with: set OPENAI_API_KEY=your-key-here")
                print("  Skipping Pass 3.\n")
            else:
                pass3_out = os.path.join(temp_dir, "output_pass3.pptx")
                print("\n" + "-" * 60)
                run_pass3(args.data, current_pptx, pass3_out)
                current_pptx = pass3_out

                if args.keep_intermediates:
                    shutil.copy2(pass3_out, os.path.join(base_dir, "output_pass3.pptx"))

        # Copy final output
        shutil.copy2(current_pptx, args.out)

        # Validate
        passes_with_sections = [p for p in passes if p == 3]
        validate_output(
            args.out,
            original_slide_count,
            num_sections if passes_with_sections else 0,
        )

    finally:
        # Clean up temp directory
        if not args.keep_intermediates:
            shutil.rmtree(temp_dir, ignore_errors=True)

    print()
    print("+" + "=" * 58 + "+")
    print(f"|   PIPELINE COMPLETE                                      |")
    print(f"|   Output: {args.out:<47}|")
    print("+" + "=" * 58 + "+")


if __name__ == "__main__":
    main()
