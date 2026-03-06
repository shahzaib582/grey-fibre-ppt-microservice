"""
PASS 2 — Add AI Restatement Sentence
=====================================
For each question slide (already processed by Pass 1):
  - Detect question number
  - Extract the numeric values from the data
  - Send to LLM for a one-sentence summary (≤35 words)
  - Replace content with the AI sentence only (no raw stats)

Usage:
  python -m survey_pipeline.pass2_add_restatement --data DATA.xlsx --pptx output_pass1.pptx --out output_pass2.pptx
"""

import argparse
import re
from pptx import Presentation
from pptx.oxml.ns import qn
import copy

from .utils import (
    PLACEHOLDER,
    parse_question_spec,
    get_question_ids,
    select_top_rows,
    select_top_rows_multi,
    format_values,
    get_slide_text,
    generate_restatement,
    set_shape_text_to_single_paragraph,
    replace_placeholder_in_shape,
    ensure_key_finding_style,
)
from .data_loader import load_ai_long


def find_values_shape(slide):
    """
    Find the shape that contains the Pass 1 values (numeric findings).
    This is the shape that previously had {Insert Finding Here} and now
    has the formatted values like 'Option – XX%; Option – XX%.'

    Also works if the placeholder hasn't been replaced yet (for standalone usage).
    """
    # Pattern: looks for "Something – NN%" which is our Pass 1 output format
    value_pattern = re.compile(r"\w+\s*[–-]\s*\d+%")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text

        # If placeholder still exists (running standalone without Pass 1)
        if PLACEHOLDER in text:
            return shape, "placeholder"

        # If Pass 1 values are present
        if value_pattern.search(text):
            return shape, "values"

    return None, None


def slide_has_table(slide):
    """True if the slide contains a table shape (e.g. multi-question favorability table)."""
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return True
    return False


def find_question_shape(slide):
    """
    Find the shape that contains the question text (e.g. 'Question 6:' line).
    Used as the style source for key-finding text.
    """
    pattern = re.compile(r"\bQuestion[s]?\s+\d+\s*:", re.IGNORECASE)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if pattern.search(text):
            return shape
    return None


def prepend_restatement_to_shape(shape, restatement: str):
    """
    Prepend the AI restatement sentence above the existing content in the shape.
    Preserves formatting from existing content.
    """
    tf = shape.text_frame

    # Get formatting reference from the first paragraph's first run
    font_name = None
    font_size = None
    font_color = None

    for para in tf.paragraphs:
        if para.runs:
            ref_run = para.runs[0]
            font_name = ref_run.font.name
            font_size = ref_run.font.size
            try:
                font_color = ref_run.font.color.rgb
            except Exception:
                pass
            break

    # Create a new paragraph element for the restatement
    first_para = tf.paragraphs[0]

    # Create new paragraph XML element
    new_p = copy.deepcopy(first_para._p)

    # Clear all runs in the new paragraph
    for r in new_p.findall(qn('a:r')):
        new_p.remove(r)

    # Create a new run with the restatement text
    new_r = copy.deepcopy(first_para.runs[0]._r) if first_para.runs else None

    if new_r is not None:
        new_r.text = restatement
        new_p.append(new_r)
    else:
        from lxml import etree
        new_r = etree.SubElement(new_p, qn('a:r'))
        rPr = etree.SubElement(new_r, qn('a:rPr'), attrib={'lang': 'en-US'})
        t = etree.SubElement(new_r, qn('a:t'))
        t.text = restatement
        new_p.append(new_r)

    # Also add an empty paragraph for spacing
    empty_p = copy.deepcopy(new_p)
    for r in empty_p.findall(qn('a:r')):
        r.text = ""

    # Insert the restatement paragraph BEFORE the first paragraph
    txBody = first_para._p.getparent()
    txBody.insert(list(txBody).index(first_para._p), new_p)

    return True


def process_slide(slide, ai_long, key_style, top_k=3, exclude_net=True):
    """Process a single slide: generate and set AI restatement (sentence only)."""
    slide_text = get_slide_text(slide)

    qspec = parse_question_spec(slide_text)
    if not qspec:
        return False

    qids = get_question_ids(qspec)
    if not qids:
        return False

    # Find the shape with values
    target_shape, mode = find_values_shape(slide)
    if target_shape is None:
        return False

    # Find the question text shape to use as style source
    question_shape = find_question_shape(slide)

    # Get data for building bullet list for LLM
    all_rows = select_top_rows_multi(ai_long, qids, top_k=top_k, exclude_net=exclude_net)
    if all_rows.empty:
        return False

    # Build bullet list for LLM prompt
    bullets = []
    for _, r in all_rows.iterrows():
        opt = str(r["answer_option"]).strip()
        pct = round(float(r["pct"]))
        bullets.append(f"- {opt} – {pct}%")
    bullet_text = "\n".join(bullets)

    # Generate restatement
    try:
        restatement = generate_restatement(bullet_text)
    except Exception as e:
        print(f"  [ERROR] LLM call failed for {qids}: {e}")
        return False

    if mode == "placeholder":
        # If placeholder still exists, replace it with restatement only (sentence, no raw stats).
        # replace_placeholder_in_shape already preserves the placeholder's font and color.
        return replace_placeholder_in_shape(target_shape, restatement)
    else:
        # Replace entire content with restatement sentence only (no raw stats),
        # and force the run to use the key-finding style we captured from the template.
        return set_shape_text_to_single_paragraph(target_shape, restatement, style=key_style)


def main():
    ap = argparse.ArgumentParser(description="Pass 2: Add AI restatement sentences")
    ap.add_argument(
        "--data",
        required=True,
        help="Path to survey Excel file (either AI-ready with 'ai_long' sheet or raw 250870-style 'ExcelData')",
    )
    ap.add_argument("--pptx", required=True, help="Path to Pass 1 output (or template for standalone)")
    ap.add_argument("--out", required=True, help="Output file path")
    ap.add_argument("--exclude-net", action="store_true", default=True, help="Exclude NET aggregates")
    ap.add_argument("--top-k", type=int, default=3, help="Number of top answers per question")
    args = ap.parse_args()

    print("=" * 60)
    print("PASS 2 — Add AI Restatement Sentences")
    print("=" * 60)
    print(f"Data:     {args.data}")
    print(f"Input:    {args.pptx}")
    print(f"Output:   {args.out}")
    print()

    ai_long = load_ai_long(args.data)
    prs = Presentation(args.pptx)

    key_style = ensure_key_finding_style(prs)

    updated = 0
    total = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        slide_text = get_slide_text(slide)
        qspec = parse_question_spec(slide_text)
        if qspec:
            qids = get_question_ids(qspec)
            if slide_has_table(slide):
                print(f"Slide {i + 1}: Skipping {qids} (slide has table — no summary sentence)")
                continue
            print(f"Slide {i + 1}: Processing {qids}...")
            if process_slide(slide, ai_long, key_style, top_k=args.top_k, exclude_net=args.exclude_net):
                print(f"  [OK] Restatement added")
                updated += 1
            else:
                print(f"  [SKIP] No values shape found or LLM failed")

    prs.save(args.out)
    print()
    print(f"PASS 2 complete.")
    print(f"  Updated: {updated} slide(s)")
    print(f"  Total:   {total} slide(s)")
    print(f"  Output:  {args.out}")


if __name__ == "__main__":
    main()
