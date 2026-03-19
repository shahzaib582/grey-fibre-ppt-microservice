"""
PASS 3 — Generate Transition Slides
====================================
For each major section (Mood, Favorability, Ballot, etc.):
  - Detect section divider slides
  - Insert exactly 2 transition slides after each divider:
    Slide A: "{Section Name}: Questions Asked" — bullet summary, no numbers
    Slide B: "{Section Name}: Survey Responses" — synthesis with key percentages
  - Content is LLM-generated but data-bound

Usage:
  python -m survey_pipeline.pass3_transition_slides --data DATA.xlsx --pptx output_pass2.pptx --out FINAL.pptx
"""

import argparse
import copy
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from lxml import etree
from pptx.oxml.ns import qn

from .utils import (
    SECTION_NAMES,
    PLACEHOLDER,
    _sanitize_pptx_text,
    is_section_divider,
    get_section_questions,
    get_section_name_for_slide,
    get_slide_text,
    get_question_text_from_slide,
    parse_question_spec,
    get_question_ids,
    generate_questions_asked_content,
    generate_survey_responses_content,
    generate_multi_question_summary_content,
    generate_executive_summary_slides,
    extract_chart_callouts_from_deck,
    call_llm,
    ensure_key_finding_style,
    apply_style_to_run,
)
from .data_loader import load_ai_long


def get_slide_layout(prs):
    """
    Get the content-style layout for transition slides (blue/white theme).
    Prefer layouts used by content slides, NOT the section divider layout.
    """
    preferred_names = [
        "Title and Content",
        "Title Only",
        "Content with Caption",
        "Blank",
        "Custom Layout",
    ]
    for layout in prs.slide_layouts:
        if layout.name in preferred_names:
            return layout
    # Use first non–section-header layout if possible (often index 1 or 2 is content)
    for idx in (1, 2, 3, 0):
        if idx < len(prs.slide_layouts):
            return prs.slide_layouts[idx]
    return prs.slide_layouts[0]


def _remap_embed_rids_in_element(element, src_part, dst_part):
    """
    Find all r:embed refs in element and add corresponding relationships
    to dst_part, updating the embed attrs to the new rIds.
    """
    r_embed = qn("r:embed")
    for el in element.iter():
        rid = el.get(r_embed)
        if rid and src_part.rels.get(rid):
            rel = src_part.rels[rid]
            new_rid = dst_part.rels.get_or_add(rel.reltype, rel.target_part)
            el.set(r_embed, new_rid)


def _copy_background_from_element(src_element, src_part, dst_slide):
    """Copy p:bg from src_element (cSld) to dst_slide. Remap image rIds to dst part."""
    src_bg = src_element.find(qn("p:bg"))
    if src_bg is None:
        return False
    dst_cSld = dst_slide._element.find(qn("p:cSld"))
    if dst_cSld is None:
        return False
    dst_bg = dst_cSld.find(qn("p:bg"))
    if dst_bg is not None:
        dst_cSld.remove(dst_bg)
    new_bg = copy.deepcopy(src_bg)
    _remap_embed_rids_in_element(new_bg, src_part, dst_slide.part)
    dst_cSld.insert(0, new_bg)
    dst_slide.follow_master_background = False
    return True


def _copy_slide_background(ref_slide, dst_slide):
    """
    Copy the background (light grey + blue footer) to dst_slide.
    Only from ref_slide directly—copying from layout can corrupt (different part/rels).
    """
    if getattr(ref_slide, "follow_master_background", True):
        return
    cSld = ref_slide._element.find(qn("p:cSld"))
    if cSld is not None:
        _copy_background_from_element(cSld, ref_slide.part, dst_slide)


def _next_shape_id(slide):
    """Return next available shape id for the slide."""
    max_id = 0
    for shp in slide.shapes:
        sid = getattr(shp, "shape_id", None)
        if sid is not None:
            try:
                max_id = max(max_id, int(sid))
            except (TypeError, ValueError):
                pass
    return max_id + 1


def _copy_slide_number_placeholder(ref_slide, dst_slide):
    """Copy the slide number placeholder from ref_slide if dst_slide doesn't have one.
    Only use ref_slide—layout/master copy can corrupt (different part/rels)."""
    has_sldnum = any(
        getattr(shp, "is_placeholder", False)
        and getattr(shp, "placeholder_format", None)
        and getattr(shp.placeholder_format, "type", None) == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER
        for shp in dst_slide.shapes
    )
    if has_sldnum or ref_slide is None:
        return
    sources = [ref_slide]  # Only ref_slide to avoid cross-part corruption
    for src in sources:
        if src is None:
            continue
        for shp in src.shapes:
            if not getattr(shp, "is_placeholder", False):
                continue
            if getattr(shp, "placeholder_format", None) is None:
                continue
            if getattr(shp.placeholder_format, "type", None) != PP_PLACEHOLDER_TYPE.SLIDE_NUMBER:
                continue
            el = shp.element
            new_el = copy.deepcopy(el)
            # Assign unique shape id to avoid PowerPoint repair errors
            nv_sp_pr = new_el.find(qn("p:nvSpPr"))
            if nv_sp_pr is not None:
                c_nv_pr = nv_sp_pr.find(qn("p:cNvPr"))
                if c_nv_pr is not None:
                    c_nv_pr.set("id", str(_next_shape_id(dst_slide)))
            # Remap r:embed refs if copying from layout (different part)
            if src is not ref_slide:
                _remap_embed_rids_in_element(new_el, src.part, dst_slide.part)
            sp_tree = dst_slide.shapes._spTree
            ext_lst = sp_tree.find(qn("p:extLst"))
            if ext_lst is not None:
                sp_tree.insert_element_before(new_el, qn("p:extLst"))
            else:
                sp_tree.append(new_el)
            return


def create_transition_slide(prs, slide_index, title_text, body_text, ref_slide=None, key_style=None, layout=None):
    """
    Create a new transition slide using a content-style layout.
    Uses the template's title and body placeholders so fonts/colors,
    background, and footer elements match the chosen layout.
    """
    title_text = _sanitize_pptx_text(title_text or "")
    body_text = _sanitize_pptx_text(body_text or "")
    if layout is None:
        layout = get_slide_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Copy the Key Findings slide's explicit background (light grey + blue footer).
    if ref_slide is not None:
        try:
            _copy_slide_background(ref_slide, slide)
        except Exception:
            pass

    # Ensure slide number placeholder is present (layout may not include it).
    if ref_slide is not None:
        try:
            _copy_slide_number_placeholder(ref_slide, slide)
        except Exception:
            pass

    # Prefer using the layout's title + body placeholders so we inherit
    # font, color, and spacing from the template. Never use slide number as title/body.
    title_shape = None
    body_shape = None

    def _is_slide_num(shp):
        return (
            getattr(shp, "is_placeholder", False)
            and getattr(shp, "placeholder_format", None)
            and getattr(shp.placeholder_format, "type", None) == PP_PLACEHOLDER_TYPE.SLIDE_NUMBER
        )

    for shape in slide.shapes:
        if not shape.has_text_frame or _is_slide_num(shape):
            continue
        if getattr(shape, "is_placeholder", False):
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE):
                title_shape = shape
            elif ph_type == PP_PLACEHOLDER_TYPE.BODY and body_shape is None:
                body_shape = shape

    # Fallbacks: first text shape as title, second as body (skip slide number)
    if title_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame and not _is_slide_num(shape):
                title_shape = shape
                break
    if body_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame and not _is_slide_num(shape) and shape is not title_shape:
                body_shape = shape
                break

    # If the chosen layout has no usable text shapes (e.g. Blank layout),
    # fall back to creating our own title/body boxes. Size them based on
    # the actual slide width so margins match other slides.
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    footer_clearance = Inches(0.85)  # leave space for blue footer so text doesn't touch
    left = Inches(0.5)
    right_margin = Inches(0.5)
    width = slide_width - left - right_margin
    title_top = Inches(0.55)
    title_height = Inches(0.5)
    body_top = Inches(1.15) + Pt(14)  # Gap after heading line
    body_height_max = slide_height - int(body_top) - int(footer_clearance)
    body_height = min(int(Inches(5.5)), body_height_max)

    if title_shape is None or not getattr(title_shape, "has_text_frame", False):
        title_shape = slide.shapes.add_textbox(left, title_top, width, title_height)
    if body_shape is None or not getattr(body_shape, "has_text_frame", False):
        body_shape = slide.shapes.add_textbox(left, body_top, width, body_height)
    else:
        # Add space after line below heading; constrain height for footer
        body_shape.top = int(body_shape.top) + int(Pt(14))
        body_max_for_layout = slide_height - int(body_shape.top) - int(footer_clearance)
        if body_shape.height > body_max_for_layout:
            body_shape.height = body_max_for_layout

    # Set title text (truncate to 9 words so it fits on one line)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        words = title_text.split()
        if len(words) > 9:
            title_text = " ".join(words[:9])
        tf = title_shape.text_frame
        tf.text = ""
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.space_after = Pt(14)  # Space after line below heading
        # Apply key-finding style to the title as well, but force bold
        if p.runs and key_style:
            apply_style_to_run(p.runs[0], key_style, force_bold=True)
            # Override size to 28pt as requested
            p.runs[0].font.size = Pt(28)

    # Set body content
    if body_shape is not None and getattr(body_shape, "has_text_frame", False):
        tf = body_shape.text_frame
        tf.text = ""
        tf.word_wrap = True
        _set_body_content(tf, body_text, key_style)

    # Clear any other shapes with placeholder text so validation finds the real title
    # (content placeholder may have "Click to add text" or {Insert Finding Here})
    for shape in slide.shapes:
        if shape is title_shape or shape is body_shape:
            continue
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if PLACEHOLDER in t or "click to add" in t.lower():
                shape.text_frame.text = ""

    return slide


# Max chars per transition slide body (excess goes to continuation slides)
MAX_CHARS_PER_SLIDE = 1150


def _build_multi_q_body_with_question_context(question_text_from_slide: str, summary_content: str) -> str:
    """Put the question text from the table slide at top of summary body (same position as on table slide)."""
    if not question_text_from_slide or not question_text_from_slide.strip():
        return summary_content
    return question_text_from_slide.strip() + "\n\n" + summary_content


def _split_body_content(body_text: str, max_chars: int = MAX_CHARS_PER_SLIDE) -> list[str]:
    """Split body text into chunks that fit on one slide. Splits at bullet boundaries."""
    lines = body_text.strip().split("\n")
    chunks = []
    current = []
    current_len = 0

    for line in lines:
        line_stripped = line.strip()
        if not line_stripped:
            continue
        line_len = len(line_stripped) + 1  # +1 for newline
        if current and current_len + line_len > max_chars:
            chunks.append("\n".join(current))
            current = []
            current_len = 0
        current.append(line_stripped)
        current_len += line_len

    if current:
        chunks.append("\n".join(current))
    return chunks if chunks else [""]


def _set_body_content(text_frame, body_text, key_style=None):
    """Set body content with proper formatting. Handles bullet points."""
    lines = body_text.strip().split("\n")

    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue

        # Determine if this is a bullet point
        is_bullet = line.startswith("•") or line.startswith("-") or line.startswith("*")
        if is_bullet:
            # Remove bullet character
            line = line.lstrip("•-* ").strip()

        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = line
        p.space_after = Pt(6)

        if is_bullet:
            p.level = 0
            # Add bullet character back for visual
            p.text = f"• {line}"

        # Apply key-finding style (same as main bullets) if available
        if p.runs and key_style:
            apply_style_to_run(p.runs[0], key_style)


def _update_table_of_contents(prs, section_names: list[str], key_style: dict | None = None):
    """Update the Table of Contents slide with actual section titles from the presentation.
    Uses key_style (same as transition slides) for font color, style, and size."""
    toc_slide = None
    body_shape = None

    for slide in prs.slides:
        full_text = get_slide_text(slide).lower()
        if "table of contents" not in full_text:
            continue

        toc_slide = slide
        # Find body shape: either a separate shape with the list, or the same shape as title (para 2+)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t:
                continue
            lines = [ln.strip() for ln in t.split("\n") if ln.strip()]
            first_line = lines[0].lower() if lines else ""
            if first_line == "table of contents":
                if len(lines) > 1:
                    # Title + list in same shape; use this shape, replace content after title
                    body_shape = shape
                    break
                continue
            # Separate body shape with the list
            body_shape = shape
            break

        if body_shape is not None:
            break

    if toc_slide is None or body_shape is None:
        return

    tf = body_shape.text_frame

    # Check if title and body are in same shape
    full_text = tf.text.strip()
    if full_text.lower().startswith("table of contents"):
        # Keep title para, replace body paras with actual section names + X placeholder for slide numbers
        paras = list(tf.paragraphs)
        txBody = tf._txBody
        for p in paras[1:]:
            txBody.remove(p._p)
        # Right-aligned tab for X column (8.5" from left)
        tab_pos_emu = int(Inches(8.5))
        for name in section_names:
            p = tf.add_paragraph()
            p.text = f"{name}\tX"
            p.level = 0
            p.space_after = Pt(6)
            # Add right-aligned tab stop so X aligns in a column
            pPr = p._p.find(qn("a:pPr"))
            if pPr is None:
                pPr = etree.SubElement(p._p, qn("a:pPr"))
            tab_lst = pPr.find(qn("a:tabLst"))
            if tab_lst is None:
                tab_lst = etree.SubElement(pPr, qn("a:tabLst"))
            etree.SubElement(tab_lst, qn("a:tab"), pos=str(tab_pos_emu), algn="r")
            if key_style and p.runs:
                apply_style_to_run(p.runs[0], key_style)
    else:
        # Body in separate shape (bullet format to match TOC style, with X placeholder)
        tf.text = ""
        tab_pos_emu = int(Inches(8.5))
        for i, name in enumerate(section_names):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {name}\tX"
            p.level = 0
            p.space_after = Pt(6)
            pPr = p._p.find(qn("a:pPr"))
            if pPr is None:
                pPr = etree.SubElement(p._p, qn("a:pPr"))
            tab_lst = pPr.find(qn("a:tabLst"))
            if tab_lst is None:
                tab_lst = etree.SubElement(pPr, qn("a:tabLst"))
            etree.SubElement(tab_lst, qn("a:tab"), pos=str(tab_pos_emu), algn="r")
            if key_style and p.runs:
                apply_style_to_run(p.runs[0], key_style)


def _replace_key_findings_with_section(prs):
    """Replace 'Key Findings' with section name on all slides in each section.
    Replaces text in-place within runs to preserve font color, family, and size.
    Also adds space after the heading line on key findings content slides (template)."""
    section_ranges = []
    for idx, slide in enumerate(prs.slides):
        sec_name = is_section_divider(slide)
        if sec_name:
            section_ranges.append((idx, sec_name))
    for j, (start_idx, section_name) in enumerate(section_ranges):
        end_idx = section_ranges[j + 1][0] if j + 1 < len(section_ranges) else len(prs.slides)
        for k in range(start_idx, end_idx):
            slide = prs.slides[k]
            title_shape = None
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if "key findings" in run.text.lower():
                            run.text = run.text.replace("Key Findings", section_name)
                            run.text = run.text.replace("key findings", section_name)
                            para.space_after = Pt(6)  # Minimal gap after line (avoid summary colliding with chart)
                            title_shape = shape
                            # If title and body are in same shape, add space_before to next para
                            paras = list(tf.paragraphs)
                            para_idx = next((i for i, p in enumerate(paras) if p is para), -1)
                            if para_idx >= 0 and para_idx + 1 < len(paras):
                                paras[para_idx + 1].space_before = Pt(6)
                            break
                    if title_shape is not None:
                        break
                if title_shape is not None:
                    break
            # Key findings content slides: title and body are separate shapes; move body down
            if title_shape is not None:
                title_bottom = int(title_shape.top) + int(title_shape.height)
                candidate = None
                for shape in slide.shapes:
                    if shape is title_shape or not getattr(shape, "has_text_frame", False):
                        continue
                    if getattr(shape, "top", None) is None:
                        continue
                    top = int(shape.top)
                    if top >= title_bottom - 5000:  # Below or slightly overlapping
                        if candidate is None or top < int(candidate.top):
                            candidate = shape
                if candidate is not None:
                    candidate.top = int(candidate.top) + int(Pt(6))  # Minimal gap (avoid summary colliding with chart)


def _find_executive_summary_slide_index(prs) -> int | None:
    """Return the index of the slide titled 'Executive Summary', or None."""
    for i, slide in enumerate(prs.slides):
        sec = is_section_divider(slide)
        if sec and "executive summary" in sec.lower():
            return i
        text = get_slide_text(slide).lower()
        if "executive summary" in text:
            lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
            if lines and lines[0].lower() == "executive summary":
                return i
    return None


def _load_exec_summary_goals(goals_path: str | None) -> str:
    """Load goals from file. Returns empty string if no path or file not found (no default)."""
    if not goals_path or not str(goals_path).strip():
        return ""
    try:
        with open(goals_path, encoding="utf-8") as f:
            return f.read().strip()
    except Exception:
        return ""


def move_slide_to_index(prs, slide, target_index):
    """Move a slide to a specific index position. Uses parent.insert to avoid XML corruption."""
    sldIdLst = prs.slides._sldIdLst
    r_id = qn("r:id")

    slide_rel_id = None
    for rel in prs.part.rels.values():
        if getattr(rel, "target_part", None) == slide.part:
            slide_rel_id = rel.rId
            break
    if not slide_rel_id:
        return

    target_entry = None
    for sldId in sldIdLst:
        if sldId.get(r_id) == slide_rel_id:
            target_entry = sldId
            break
    if target_entry is None:
        return

    sldIdLst.remove(target_entry)
    entries = list(sldIdLst)
    if target_index >= len(entries):
        sldIdLst.append(target_entry)
    else:
        # Use parent.insert to avoid addprevious corruption
        parent = target_entry.getparent() or sldIdLst
        idx = list(parent).index(entries[target_index])
        parent.insert(idx, target_entry)


def main():
    ap = argparse.ArgumentParser(description="Pass 3: Generate transition slides")
    ap.add_argument(
        "--data",
        required=True,
        help="Path to survey Excel file (either AI-ready with 'ai_long' sheet or raw 250870-style 'ExcelData')",
    )
    ap.add_argument("--pptx", required=True, help="Path to Pass 2 output")
    ap.add_argument("--out", required=True, help="Output file path")
    ap.add_argument(
        "--exec-summary-goals",
        default=None,
        help="Path to file with survey goals for executive summary (required; no default)",
    )
    args = ap.parse_args()

    print("=" * 60)
    print("PASS 3 — Generate Transition Slides")
    print("=" * 60)
    print(f"Data:     {args.data}")
    print(f"Input:    {args.pptx}")
    print(f"Output:   {args.out}")
    print()

    ai_long = load_ai_long(args.data)
    prs = Presentation(args.pptx)

    key_style = ensure_key_finding_style(prs)

    original_count = len(prs.slides)
    print(f"Original slide count: {original_count}")
    print()

    # Step 1: Find all section dividers and their data
    sections_found = []
    slides_list = list(prs.slides)

    for i, slide in enumerate(slides_list):
        section_name = is_section_divider(slide)
        if section_name:
            section_data = get_section_questions(prs, i, ai_long)
            sections_found.append({
                "name": section_name,
                "slide_index": i,
                "slide": slide,
                "questions": section_data,
            })
            print(f"Found section: '{section_name}' at slide {i + 1}")
            print(f"  Questions: {section_data['question_ids']}")

    print(f"\nTotal sections found: {len(sections_found)}")
    print()

    if not sections_found:
        print("[WARN] No section dividers found. Nothing to do.")
        prs.save(args.out)
        return

    # Update Table of Contents slide with actual section titles (use transition slide style)
    section_names = [s["name"] for s in sections_found]
    _update_table_of_contents(prs, section_names, key_style=key_style)
    print(f"Updated Table of Contents with {len(section_names)} section(s)")

    # Use the actual Key Findings CONTENT slide (light grey + blue footer).
    # Exclude the dark-blue "Key Findings" divider. Content slides have a
    # chart and/or substantial body text; the divider has little content.
    keyfinding_layout = None
    keyfinding_slide = None
    candidates = []
    for s in prs.slides:
        text = get_slide_text(s)
        if "key findings" not in text.lower():
            continue
        if is_section_divider(s):
            continue
        has_chart = any(getattr(shp, "has_chart", False) for shp in s.shapes)
        candidates.append((s, len(text), has_chart))
    # Prefer: has chart, then longest text (content has bullet + chart labels)
    candidates.sort(key=lambda x: (x[2], x[1]), reverse=True)
    if candidates:
        keyfinding_slide, _, _ = candidates[0]
        keyfinding_layout = keyfinding_slide.slide_layout or get_slide_layout(prs)
        print(f"Using Key Findings content slide for transitions: '{keyfinding_layout.name}'")
    if keyfinding_layout is None:
        keyfinding_layout = get_slide_layout(prs)
        print(f"[INFO] No Key Findings content slide found; falling back to '{keyfinding_layout.name}'")

    slides_inserted = 0
    # Pre-generate transition content and multi-q content (reused for exec summary + slide insertion)
    print("Pre-generating transition and multi-question content...")
    section_transition_content = {}  # section_name -> (questions_asked, survey_responses)
    for section in sections_found:
        if section["name"] == "Executive Summary" or not section["questions"]["question_ids"]:
            continue
        q_data = section["questions"]
        qa, sr = "• Questions in this section.", "• Survey response data available."
        try:
            qa = generate_questions_asked_content(section["name"], q_data["question_texts"])
        except Exception:
            pass
        try:
            sr = generate_survey_responses_content(
                section["name"], q_data["question_texts"], q_data["question_data"]
            )
        except Exception:
            pass
        section_transition_content[section["name"]] = (qa, sr)

    multi_q_content = []  # list of (slide_idx, section_name, question_text_from_slide, question_data, summary)
    for i, slide in enumerate(prs.slides):
        if is_section_divider(slide):
            continue
        slide_text = get_slide_text(slide)
        qspec = parse_question_spec(slide_text)
        if not qspec or (qspec[0] != "range" and len(get_question_ids(qspec)) <= 1):
            continue
        qids = get_question_ids(qspec)
        section_name = get_section_name_for_slide(prs, i)
        question_text_from_slide = get_question_text_from_slide(slide)
        question_texts = {}
        for qid in qids:
            rows = ai_long[ai_long["question_id"] == qid]
            if not rows.empty:
                qt = rows.iloc[0].get("question_text", qid) if "question_text" in rows.columns else qid
                question_texts[qid] = str(qt).strip() if pd.notna(qt) and str(qt).strip() else qid
        question_data = ai_long[ai_long["question_id"].isin(qids)].copy()
        mq_summary = "• Summary of questions and findings."
        try:
            mq_summary = generate_multi_question_summary_content(
                section_name, question_texts, question_data
            )
        except Exception:
            pass
        multi_q_content.append((i, section_name, question_text_from_slide, question_data, mq_summary))

    # Step 2: Insert transition slides FIRST (work backwards) so exec summary block stays together
    transition_counts = {}  # section_name -> number of slides inserted
    for section in reversed(sections_found):
        section_name = section["name"]
        q_data = section["questions"]
        divider_idx = section["slide_index"]

        if not q_data["question_ids"]:
            print(f"[SKIP] '{section_name}' — no questions found")
            continue

        print(f"Inserting transition slides for '{section_name}'...")
        questions_content, responses_content = section_transition_content.get(
            section_name, ("• Questions were asked in this section.", "• Survey response data available for review.")
        )

        # Create Slide B — inserted first so it ends up AFTER Slide A (title = section name only)
        # Split content; excess goes to continuation slides. Insert in reverse so order is correct.
        response_chunks = _split_body_content(responses_content)
        for chunk in reversed(response_chunks):
            slide_b = create_transition_slide(
                prs,
                divider_idx + 1,
                section_name,
                chunk,
                ref_slide=keyfinding_slide,
                key_style=key_style,
                layout=keyfinding_layout,
            )
            move_slide_to_index(prs, slide_b, divider_idx + 1)
            slides_inserted += 1

        # Create Slide A (title = section name only)
        question_chunks = _split_body_content(questions_content)
        for chunk in reversed(question_chunks):
            slide_a = create_transition_slide(
                prs,
                divider_idx + 1,
                section_name,
                chunk,
                ref_slide=keyfinding_slide,
                key_style=key_style,
                layout=keyfinding_layout,
            )
            move_slide_to_index(prs, slide_a, divider_idx + 1)
            slides_inserted += 1

        total = len(response_chunks) + len(question_chunks)
        transition_counts[section_name] = total
        print(f"  [OK] {total} transition slide(s) inserted after '{section_name}'")

    # Step 2.5: Insert executive summary slides (after transition slides so indices are stable)
    exec_summary_idx = _find_executive_summary_slide_index(prs)
    orig_exec_summary_idx = next(
        (s["slide_index"] for s in sections_found if "executive summary" in s["name"].lower()),
        None,
    )
    num_exec_slides = 0
    goals_path = getattr(args, "exec_summary_goals", None)
    goals_text = _load_exec_summary_goals(goals_path)
    if exec_summary_idx is not None and not goals_text:
        print("[INFO] Executive summary skipped: --exec-summary-goals file required (no default)")
    if exec_summary_idx is not None and goals_text:
        sections_for_exec = [s for s in sections_found if s["name"] != "Executive Summary"]
        summary_parts = []
        callouts = extract_chart_callouts_from_deck(prs, ai_long)
        if callouts:
            by_sec = {}
            for c in callouts:
                sec = c["section"]
                if sec not in by_sec:
                    by_sec[sec] = []
                by_sec[sec].append(f"  - {c['question_text'][:80]}... → {c['summary']}")
            for sec, items in by_sec.items():
                summary_parts.append(f"Chart callouts ({sec}):\n" + "\n".join(items))
        for sec_name, (qa, sr) in section_transition_content.items():
            summary_parts.append(f"Questions Asked ({sec_name}):\n{qa}")
            summary_parts.append(f"Survey Responses ({sec_name}):\n{sr}")
        for _, sec_name, _, _, mq in multi_q_content:
            summary_parts.append(f"Multi-question summary ({sec_name}):\n{mq}")
        summaries_text = "\n\n".join(summary_parts) if summary_parts else None
        try:
            exec_slides_data = generate_executive_summary_slides(
                goals_text, sections_for_exec, ai_long, summaries_text=summaries_text
            )
            if exec_slides_data:
                num_exec_slides = len(exec_slides_data)
                for slide_data in reversed(exec_slides_data):
                    body_text = "\n".join(f"• {b}" for b in slide_data.get("bullets", []))
                    if not body_text.strip():
                        body_text = "• Key findings from the survey."
                    exec_slide = create_transition_slide(
                        prs,
                        exec_summary_idx + 1,
                        slide_data.get("title", "Key Finding"),
                        body_text,
                        ref_slide=keyfinding_slide,
                        key_style=key_style,
                        layout=keyfinding_layout,
                    )
                    move_slide_to_index(prs, exec_slide, exec_summary_idx + 1)
                    slides_inserted += 1
                print(f"Inserted {num_exec_slides} executive summary slide(s) after slide {exec_summary_idx + 1}")
        except Exception as e:
            print(f"[ERROR] Executive summary generation failed: {e}")

    # Step 2.6: Insert multi-question summary slides BEFORE the question/table slide (work backwards)
    # Adjust slide_idx for transition slides and exec slides inserted before each multi-q slide
    def _adjusted_multi_q_idx(orig_idx: int) -> int:
        adj = orig_idx
        for s in sections_found:
            if s["name"] == "Executive Summary" or not s["questions"]["question_ids"]:
                continue
            if s["slide_index"] + 1 <= orig_idx:
                adj += transition_counts.get(s["name"], 0)
        if orig_exec_summary_idx is not None and orig_idx > orig_exec_summary_idx:
            adj += num_exec_slides
        return adj

    for slide_idx, section_name, question_text_from_slide, _qd, summary_content in reversed(multi_q_content):
        adj_idx = _adjusted_multi_q_idx(slide_idx)
        body_with_context = _build_multi_q_body_with_question_context(question_text_from_slide, summary_content)
        slide_summary = create_transition_slide(
            prs,
            adj_idx,
            section_name,
            body_with_context,
            ref_slide=keyfinding_slide,
            key_style=key_style,
            layout=keyfinding_layout,
        )
        move_slide_to_index(prs, slide_summary, adj_idx)
        slides_inserted += 1
        print(f"  [OK] Multi-question summary slide inserted before slide {adj_idx + 1}")

    # Step 3: Replace "Key Findings" with section name on all slides in each section
    _replace_key_findings_with_section(prs)

    # Step 3.5: Skip normalize_slide_numbers—copying placeholders between slides can corrupt PPTX.
    # Layout-provided placeholders are left as-is; new slides without placeholders stay without.

    # Step 4: Save
    prs.save(args.out)

    final_count = len(prs.slides)
    print()
    print(f"PASS 3 complete.")
    print(f"  Sections processed: {len(sections_found)}")
    print(f"  Slides inserted:    {slides_inserted}")
    print(f"  Original count:     {original_count}")
    print(f"  Final count:        {final_count}")
    print(f"  Expected count:     {original_count + slides_inserted}")
    print(f"  Output:             {args.out}")

    # Validation
    if final_count != original_count + slides_inserted:
        print(f"\n  [WARN] Slide count mismatch! Expected {original_count + slides_inserted}, got {final_count}")
    else:
        print(f"\n  [OK] Slide count validated [OK]")


if __name__ == "__main__":
    main()
